from pptx import Presentation
from pptx.util import Inches, Pt

# Повторная инициализация после сброса
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Имя автора
author_name = "Данила Пеляев"

# Слайд 1: Титульный
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Современные средства связи. Сотовая связь"
slide.placeholders[1].text = f"Автор: {author_name}\nКласс: 8\nДата: Май 2025"

# Слайд 2: Введение
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Введение"
slide.placeholders[1].text = (
    "Связь — основа современного общества.\n"
    "Сотовая связь — ключевой инструмент общения в XXI веке.\n"
    "Цель: Рассмотреть развитие, устройство и перспективы сотовой связи."
)

# Слайд 3: Виды современных средств связи (таблица)
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Виды современных средств связи"
table = slide.shapes.add_table(5, 3, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)).table
headings = ["Тип связи", "Пример", "Особенности"]
data = [
    ["Проводная", "Телефон, интернет (ADSL)", "Стабильная, но не мобильная"],
    ["Беспроводная", "Wi-Fi, Bluetooth", "Удобство, малая дальность"],
    ["Спутниковая", "Starlink, GPS", "Дорогая, глобальное покрытие"],
    ["Сотовая", "4G, 5G", "Высокая скорость и мобильность"],
]
for i in range(3):
    table.cell(0, i).text = headings[i]
for r, row_data in enumerate(data, start=1):
    for c, val in enumerate(row_data):
        table.cell(r, c).text = val

# Слайд 4: История сотовой связи
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "История сотовой связи"
slide.placeholders[1].text = (
    "1G (1980-е) — аналоговая передача.\n"
    "2G (1990-е) — цифровой сигнал, SMS.\n"
    "3G — мобильный интернет.\n"
    "4G — высокая скорость.\n"
    "5G — минимальная задержка, IoT."
)

# Слайд 5: Как работает сотовая связь
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Как работает сотовая связь"
slide.placeholders[1].text = (
    "1. Телефон отправляет сигнал на базовую станцию.\n"
    "2. Сигнал передаётся через сеть оператора.\n"
    "3. Обмен данными между устройствами.\n"
    "4. Ответный сигнал поступает обратно."
)

# Слайд 6: Сравнение поколений связи
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Сравнение поколений связи"
table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)).table
headings = ["Поколение", "Скорость", "Технология", "Особенности"]
data = [
    ["2G", "до 0.2 Мбит/с", "GSM", "Голос, SMS"],
    ["3G", "до 2 Мбит/с", "UMTS", "Видеосвязь, интернет"],
    ["4G", "до 100 Мбит/с", "LTE", "Стриминг, онлайн-игры"],
    ["5G", "до 10 Гбит/с", "mmWave", "Умные города, IoT"],
]
for i in range(4):
    table.cell(0, i).text = headings[i]
for r, row_data in enumerate(data, start=1):
    for c, val in enumerate(row_data):
        table.cell(r, c).text = val

# Слайд 7: Применение сотовой связи
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Применение сотовой связи"
slide.placeholders[1].text = (
    "• Общение (звонки, мессенджеры)\n"
    "• Интернет и социальные сети\n"
    "• Умные устройства (IoT)\n"
    "• Навигация и транспорт\n"
    "• Медицина и образование"
)

# Слайд 8: Преимущества и недостатки
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Преимущества и недостатки"
slide.placeholders[1].text = (
    "Преимущества:\n"
    "• Мобильность\n"
    "• Высокая скорость\n"
    "• Глобальное покрытие\n\n"
    "Недостатки:\n"
    "• Зависимость от сигнала\n"
    "• Возможный вред здоровью\n"
    "• Высокая нагрузка на сеть"
)

# Слайд 9: Будущее сотовой связи
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Будущее сотовой связи"
slide.placeholders[1].text = (
    "• Развитие 6G\n"
    "• Виртуальная и дополненная реальность\n"
    "• Автоматизация транспорта\n"
    "• Более надёжные и защищённые сети"
)

# Слайд 10: Заключение
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Заключение"
slide.placeholders[1].text = (
    "Сотовая связь — ключевой элемент цифрового мира.\n"
    "Она продолжает развиваться и менять нашу жизнь.\n"
    "Мы зависим от неё, но и получаем огромные преимущества."
)

# Сохраняем файл
file_path = "/mnt/data/Современные_средства_связи_Данила_Пеляев.pptx"
prs.save(file_path)
file_path
